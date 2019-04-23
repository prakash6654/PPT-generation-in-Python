#Author:Prakash

import pandas as pd
import numpy as np

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.chart import XL_TICK_MARK
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.util import Pt
from pptx.enum.chart import XL_LEGEND_POSITION
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.enum.chart import XL_LABEL_POSITION

#Data Import
raw_data = pd.read_excel('C:\Users\Prakash\Documents\InputData.xlsx' ,'Sheet1')
#Filtering required columns
data=raw_data[['Number','Created','Priority', 'Month']]

#Slicing the data
hour_data=(raw_data['Created'].dt.strftime('%H'))
monthly_data=(raw_data['Created'].dt.strftime('%b-%y'))
week_data=(raw_data['Created'].dt.strftime('%a'))

#Working with PPTX library to add slide and giving title
prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[5])
title = slide.shapes.title
title.text = "Ticket Arrival Trend"
title.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER


#1.PRIORITY CHART


#Grouping data based on priority
by_priority=raw_data.groupby(['Priority']).count()
# Get priority value to show in legend
priority=sorted(raw_data['Priority'].unique())
#Calculating the values to plot
percentage=np.array(by_priority['Number']/by_priority['Number'].sum())

# Working with Chart
chart_data = ChartData()
#Assigning Chart labels & Chart Values
chart_data.categories = priority
chart_data.add_series('Priority Based', percentage)

#Possitioning the Chart in slide
x, y, cx, cy = Inches(5.5), Inches(4.3), Inches(4.5), Inches(3.2)
#Adding Chart in Slide with title
priority_chart = slide.shapes.add_chart(XL_CHART_TYPE.PIE, x, y, cx, cy, chart_data).chart
priority_chart.chart_title.text_frame.text='By Priority'

#Chart Formatting
priority_chart.has_legend = True
priority_chart.legend.include_in_layout = False
priority_chart.legend.position = XL_LEGEND_POSITION.BOTTOM
priority_chart.legend.font.size = Pt(9)
priority_chart.plots[0].has_data_labels = True
data_labels = priority_chart.plots[0].data_labels
data_labels.number_format = '0%'
data_labels.font.color.rgb = RGBColor(0,0,0)
data_labels.font.size = Pt(12)


#MONTHLY TICKET ARRIVAL



#Replacing Created field with month value
data.loc[:,('Created')]=monthly_data

#Grouping data based on Month
by_month=data.groupby(['Created']).count()
# Get priority value to show in legend
label=sorted(data['Created'].unique() )

label=[str(a) for a in label]
month=np.array(by_month['Number'])

chart_data = ChartData()
chart_data.categories = label
chart_data.add_series('Monthly arrival', month)


x, y, cx, cy = Inches(0.2), Inches(1.2), Inches(5), Inches(3.2)
monthly_chart = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS_STACKED, x, y, cx, cy, chart_data).chart
monthly_chart.chart_title.text_frame.text='Monthly Ticket Arrival'
monthly_chart.has_legend = False
monthly_chart.plots[0].has_data_labels = True
data_labels = monthly_chart.plots[0].data_labels
data_labels.font.size = Pt(12)
category_axis = monthly_chart.category_axis
category_axis.has_major_gridlines = False
category_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
category_axis.tick_labels.font.italic = True
category_axis.tick_labels.font.size = Pt(9)

value_axis = monthly_chart.value_axis
value_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
value_axis.has_minor_gridlines = False

tick_labels = value_axis.tick_labels
tick_labels.font.size = Pt(9)



#HOURLY TICKET ARRIVAL



print(data)
data.loc[:,('Created')]=hour_data

by_hour=data.groupby(['Created']).count()

label=sorted(data['Created'].unique() )
label=[str(a) for a in label]

hour=np.array(by_hour['Number'])

chart_data = ChartData()
chart_data.categories = label
chart_data.add_series('Hourly arrival', hour)


x, y, cx, cy = Inches(0.2), Inches(4.3), Inches(5), Inches(3.2)
hourly_chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart
hourly_chart.chart_title.text_frame.text='Hourly Ticket Arrival'
plot = hourly_chart.plots[0]
plot.has_data_labels = True
data_labels = plot.data_labels
data_labels.number_format = '0'
data_labels.font.size = Pt(12)
data_labels.font.color.rgb = RGBColor(0, 0, 0)
data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
category_axis = hourly_chart.category_axis
category_axis.has_major_gridlines = False
category_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
category_axis.tick_labels.font.italic = True
category_axis.tick_labels.font.size = Pt(9)

value_axis = hourly_chart.value_axis
value_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
value_axis.has_minor_gridlines = False

tick_labels = value_axis.tick_labels
tick_labels.font.size = Pt(9)



#WEEKLY TICKET ARRIVAL


data.loc[:,('Created')]=week_data
by_week=data.groupby(['Created']).count()
label=sorted(data['Created'].unique() )


week=np.array(by_week['Number']/by_week['Number'].sum())

chart_data = ChartData()
chart_data.categories = label
chart_data.add_series('Series 1', week)


x, y, cx, cy = Inches(5.5), Inches(1.2), Inches(4.5), Inches(3.2)
chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, x, y, cx, cy, chart_data).chart
chart.chart_title.text_frame.text='Weekday Ticket Arrival'
plot = chart.plots[0]
plot.has_data_labels = True
data_labels = plot.data_labels
data_labels.number_format = '0%'
data_labels.font.size = Pt(12)
data_labels.font.color.rgb = RGBColor(0, 0, 0)
data_labels.position = XL_LABEL_POSITION.OUTSIDE_END

category_axis = chart.category_axis
category_axis.has_major_gridlines = False
category_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
category_axis.tick_labels.font.italic = True
category_axis.tick_labels.font.size = Pt(9)

value_axis = chart.value_axis
value_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
value_axis.has_minor_gridlines = False
value_axis.number_format='0%'
tick_labels = value_axis.tick_labels
tick_labels.font.size = Pt(9)
tick_labels.number_format = '0%'


#Saving the presentation
prs.save('TECH0001.pptx')
